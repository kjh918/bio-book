"""
gencurix_report.slide_utils
=============================

Low level helpers that python-pptx does not provide out of the box:
- duplicating a slide (with its images/relationships intact)
- moving a slide to a new position
- deleting a slide
- finding a shape (recursively through groups) by its placeholder text
- writing text / bullet lists / tables / images into a shape while keeping
  the template's existing formatting wherever possible
"""

import copy
from typing import Iterator, List, Optional, Tuple

from pptx.oxml.ns import qn
from pptx.util import Emu
from pptx.shapes.base import BaseShape


# ---------------------------------------------------------------------------
# Slide-level operations
# ---------------------------------------------------------------------------

def duplicate_slide(prs, index: int):
    """Duplicate the slide at `index` (0-based). The copy is appended at the
    end of the deck (use move_slide() afterwards to place it). Returns the
    new Slide object.

    Handles picture / media relationships so images copy correctly, which
    python-pptx does not do if you simply deepcopy the XML.
    """
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)

    # The new slide comes with auto-generated placeholder shapes (from the
    # layout) -- wipe them out, we're replacing spTree wholesale.
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    # Copy every shape element from the source slide's spTree.
    for shp in source.shapes:
        new_el = copy.deepcopy(shp._element)
        dest.shapes._spTree.append(new_el)
        _rewrite_relationship_ids(shp._element, new_el, source.part, dest.part)

    # Copy slide-level background / formatting overrides if present.
    src_bg = source._element.find(qn('p:cSld') + '/' + qn('p:bg'))
    if src_bg is not None:
        dest_cSld = dest._element.find(qn('p:cSld'))
        dest_bg_existing = dest_cSld.find(qn('p:bg'))
        if dest_bg_existing is not None:
            dest_cSld.remove(dest_bg_existing)
        dest_cSld.insert(0, copy.deepcopy(src_bg))

    return dest


def _rewrite_relationship_ids(src_el, dest_el, src_part, dest_part):
    """Any r:embed / r:id / r:link attribute in a copied shape points at a
    relationship ID that only exists in the *source* slide part. Walk the
    copied XML, and for every such attribute, create the matching
    relationship on the destination part (re-using the same target part --
    e.g. the same image bytes) and rewrite the attribute to the new rId.
    """
    r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    attrs_to_check = [qn('r:embed'), qn('r:id'), qn('r:link')]

    src_nodes = list(src_el.iter())
    dest_nodes = list(dest_el.iter())
    for src_node, dest_node in zip(src_nodes, dest_nodes):
        for attr in attrs_to_check:
            rid = src_node.get(attr)
            if not rid:
                continue
            try:
                rel = src_part.rels[rid]
            except KeyError:
                continue
            if rel.is_external:
                new_rid = dest_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
            else:
                new_rid = dest_part.relate_to(rel.target_part, rel.reltype)
            dest_node.set(attr, new_rid)


def move_slide(prs, old_index: int, new_index: int):
    """Move the slide currently at old_index to new_index (0-based), shifting
    the others accordingly. Operates directly on <p:sldIdLst>."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])


def delete_slide(prs, index: int):
    """Remove a slide entirely (drops it from sldIdLst; does not attempt to
    scrub now-orphaned media parts -- harmless bloat, PowerPoint ignores it)."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])


# --- ID-based variants -----------------------------------------------------
# Slide *indices* shift every time a slide is inserted/removed/reordered, so
# code that grabs several slide references up front and then mutates the
# deck should use these instead -- they key off each slide's stable
# `slide_id`, which never changes for the lifetime of the Slide object.

def _find_sldId_element(prs, slide):
    for sldId in prs.slides._sldIdLst:
        if int(sldId.get('id')) == slide.slide_id:
            return sldId
    raise ValueError(f"slide_id {slide.slide_id} not found in sldIdLst")


def delete_slide_obj(prs, slide):
    """Remove `slide` from the deck, regardless of its current index."""
    lst = prs.slides._sldIdLst
    lst.remove(_find_sldId_element(prs, slide))


def move_slide_after(prs, slide, after_slide):
    """Reposition `slide` so it comes immediately after `after_slide`."""
    lst = prs.slides._sldIdLst
    el = _find_sldId_element(prs, slide)
    lst.remove(el)
    after_el = _find_sldId_element(prs, after_slide)
    after_index = list(lst).index(after_el)
    lst.insert(after_index + 1, el)


def move_slide_before(prs, slide, before_slide):
    """Reposition `slide` so it comes immediately before `before_slide`."""
    lst = prs.slides._sldIdLst
    el = _find_sldId_element(prs, slide)
    lst.remove(el)
    before_el = _find_sldId_element(prs, before_slide)
    before_index = list(lst).index(before_el)
    lst.insert(before_index, el)


# ---------------------------------------------------------------------------
# Shape discovery
# ---------------------------------------------------------------------------

def iter_shapes_recursive(shapes) -> Iterator[BaseShape]:
    """Yield every shape in `shapes`, descending into groups."""
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:  # GROUP
            yield from iter_shapes_recursive(shape.shapes)


def find_shape_by_text(slide, marker: str) -> Optional[BaseShape]:
    """Return the first shape (searched recursively through groups) whose
    text frame's text exactly equals `marker` (after stripping whitespace)."""
    for shape in iter_shapes_recursive(slide.shapes):
        if shape.has_text_frame and shape.text_frame.text.strip() == marker:
            return shape
    return None


def find_shapes_starting_with(slide, prefix: str) -> List[BaseShape]:
    """Return all shapes (recursively) whose text starts with `prefix`."""
    out = []
    for shape in iter_shapes_recursive(slide.shapes):
        if shape.has_text_frame and shape.text_frame.text.strip().startswith(prefix):
            out.append(shape)
    return out


def find_group_containing(slide, marker: str):
    """Return the *innermost* group shape whose direct children include a
    text shape matching `marker` exactly (searched recursively through
    nested groups). Used for the TOC quadrants and Summary quadrants, where
    each quadrant is its own group holding a label + a content textbox."""

    def _search(shapes):
        for shape in shapes:
            if shape.shape_type == 6:  # GROUP
                for child in shape.shapes:
                    if child.has_text_frame and child.text_frame.text.strip() == marker:
                        return shape
                found = _search(shape.shapes)
                if found is not None:
                    return found
        return None

    return _search(slide.shapes)


# ---------------------------------------------------------------------------
# Text writing (keeps the template's font/size/color by editing the existing
# run(s) in place rather than creating brand new ones)
# ---------------------------------------------------------------------------

def set_single_line_text(shape: BaseShape, text: str):
    """Replace the text of a shape that has exactly one meaningful run,
    preserving that run's formatting."""
    tf = shape.text_frame
    first_para = tf.paragraphs[0]
    if not first_para.runs:
        first_para.text = text
        return
    first_para.runs[0].text = text
    # Remove any extra runs in the first paragraph beyond the first.
    for extra in first_para.runs[1:]:
        extra._r.getparent().remove(extra._r)
    # Remove any extra paragraphs beyond the first.
    for extra_p in tf.paragraphs[1:]:
        extra_p._p.getparent().remove(extra_p._p)


def fill_bullets(shape: BaseShape, items: List[str], bulleted: bool = True):
    """Replace a text box's content with one paragraph per string in
    `items`, cloning the box's first paragraph (and its bullet/indent/font
    formatting) for every line so the template's look is preserved.
    If `bulleted` is False, bullet formatting is left untouched but the
    template is expected to already be a non-bulleted paragraph style."""
    tf = shape.text_frame
    paragraphs = tf.paragraphs
    template_p = paragraphs[0]._p
    parent = template_p.getparent()
    template_p_copy = copy.deepcopy(template_p)

    # Remove all existing paragraphs.
    for p in list(tf.paragraphs):
        parent.remove(p._p)

    if not items:
        items = [""]

    for i, text in enumerate(items):
        new_p = copy.deepcopy(template_p_copy)
        # Keep only the first run's formatting; drop extra runs, set text.
        runs = new_p.findall(qn('a:r'))
        if runs:
            first_r = runs[0]
            for extra in runs[1:]:
                new_p.remove(extra)
            t_el = first_r.find(qn('a:t'))
            if t_el is None:
                t_el = first_r.makeelement(qn('a:t'), {})
                first_r.append(t_el)
            t_el.text = text
        else:
            # no runs at all -- just set text via paragraph text (rare)
            pass
        parent.append(new_p)


def replace_text_token(shape: BaseShape, token: str, replacement: str, paragraph_index: int = 0) -> bool:
    """Replace `token` (e.g. "[TITLE]") with `replacement` inside a
    paragraph, even when the token's characters are split across multiple
    <a:r> runs (common in PowerPoint-authored files, e.g. "Results : [" +
    "TITLE]"). The run whose span contains the token's formatting is used
    for the resulting merged run. Returns True if the token was found and
    replaced."""
    tf = shape.text_frame
    if paragraph_index >= len(tf.paragraphs):
        return False
    para = tf.paragraphs[paragraph_index]
    runs = para.runs
    if not runs:
        return False
    full_text = "".join(r.text for r in runs)
    if token not in full_text:
        return False
    new_full = full_text.replace(token, replacement)

    first_r = runs[0]._r
    format_source_r = runs[-1]._r  # formatting of the last run covering the token tends to be the "content" style
    src_rPr = format_source_r.find(qn('a:rPr'))

    t_el = first_r.find(qn('a:t'))
    if t_el is None:
        t_el = first_r.makeelement(qn('a:t'), {})
        first_r.append(t_el)
    t_el.text = new_full

    if src_rPr is not None:
        existing_rPr = first_r.find(qn('a:rPr'))
        new_rPr = copy.deepcopy(src_rPr)
        if existing_rPr is not None:
            first_r.replace(existing_rPr, new_rPr)
        else:
            first_r.insert(0, new_rPr)

    for extra in runs[1:]:
        extra._r.getparent().remove(extra._r)
    return True


def fill_table(slide, shape: BaseShape, headers: List[str], rows: List[List[str]],
                box: Optional[Tuple[int, int, int, int]] = None):
    """Remove the placeholder text-box `shape` and add a native PowerPoint
    table in its place. By default uses `shape`'s own position/size; pass
    `box=(left, top, width, height)` (EMU) to pin an explicit position/size
    instead. Returns the new table shape (GraphicFrame)."""
    if box is not None:
        left, top, width, height = box
    else:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
    shape._element.getparent().remove(shape._element)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    graphic_frame = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = graphic_frame.table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = str(h)
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.size = run.font.size or None

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            table.cell(r, c).text = str(val)

    return graphic_frame


def fill_image(slide, shape: BaseShape, image_source,
                box: Optional[Tuple[int, int, int, int]] = None):
    """Remove the placeholder shape and insert a picture, scaled to fit
    inside a bounding box (preserving aspect ratio, centered). By default
    the box is `shape`'s own position/size; pass `box=(left, top, width,
    height)` (EMU) to pin an explicit position/size instead.

    `image_source` may be a file path (str) OR a file-like object (e.g. a
    BytesIO from render_plot_to_stream()) -- both are accepted transparently
    by PIL and python-pptx. Returns the new picture shape."""
    from PIL import Image as PILImage

    if box is not None:
        left, top, box_w, box_h = box
    else:
        left, top, box_w, box_h = shape.left, shape.top, shape.width, shape.height
    shape._element.getparent().remove(shape._element)

    with PILImage.open(image_source) as im:
        img_w, img_h = im.size
    if hasattr(image_source, "seek"):
        image_source.seek(0)  # PIL consumed the stream reading the size above

    img_ratio = img_w / img_h
    box_ratio = box_w / box_h

    if img_ratio > box_ratio:
        new_w = box_w
        new_h = int(box_w / img_ratio)
    else:
        new_h = box_h
        new_w = int(box_h * img_ratio)

    new_left = left + (box_w - new_w) // 2
    new_top = top + (box_h - new_h) // 2

    return slide.shapes.add_picture(image_source, new_left, new_top, new_w, new_h)


def render_plot_to_stream(figure, dpi: int = 200):
    """Render a matplotlib Figure (or Axes -- its .figure is used) to an
    in-memory PNG stream with a sensible default export style: white
    background and a tight bounding box, so callers don't need to call
    fig.savefig() themselves. Returns a BytesIO positioned at 0."""
    import io
    fig = figure.figure if hasattr(figure, "figure") and not hasattr(figure, "savefig") else figure
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight", facecolor="white")
    buf.seek(0)
    return buf


def dataframe_to_table(df, float_format: str = "{:.3g}"):
    """Convert a pandas DataFrame into (headers, rows) of strings, ready for
    fill_table(). Floats are formatted with `float_format` by default;
    everything else is just str()'d."""
    headers = [str(c) for c in df.columns]
    rows = []
    for _, record in df.iterrows():
        formatted = []
        for value in record:
            if isinstance(value, float):
                formatted.append(float_format.format(value))
            else:
                formatted.append(str(value))
        rows.append(formatted)
    return headers, rows
