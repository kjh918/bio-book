"""
gencurix_report.builder
==========================

ReportBuilder: takes a `ReportData` object (see models.py) and the Gencurix
template, and produces a filled .pptx.

    from gencurix_report import ReportBuilder
    from gencurix_report.models import ReportData, ...

    data = ReportData(...)
    ReportBuilder().build(data, "output.pptx")
"""

import os
from pptx import Presentation

from . import slide_utils as su
from .fill_content import apply_content_block
from .models import ReportData, DividerData, TitlePageData, SummaryData

_HERE = os.path.dirname(__file__)
DEFAULT_TEMPLATE = os.path.join(_HERE, "assets", "Gencurix_PPT_Template.pptx")

# 0-based slide indices in the *original, unmodified* template. Only used at
# the very start of build(), before any slide is inserted/removed/moved --
# from that point on we hold direct Slide object references instead, since
# indices drift as soon as the deck is mutated.
IDX_COLOR_GUIDE = 0
IDX_TITLE = 1
IDX_TOC = 2
IDX_DIVIDER_1 = 3
IDX_SUMMARY = 4
IDX_INTRO = 5
IDX_WORKFLOW = 6
IDX_SAMPLE = 7
IDX_DIVIDER_2 = 8
IDX_RESULTS = 9
IDX_DIVIDER_3 = 10
IDX_CONCLUSION = 11
IDX_END = 12

TOC_MAX_CHAPTERS = 4


class ReportBuilder:
    def __init__(self, template_path: str = None):
        self.template_path = template_path or DEFAULT_TEMPLATE
        self.prs = Presentation(self.template_path)

    # ------------------------------------------------------------------
    def build(self, data: ReportData, output_path: str) -> str:
        slides = list(self.prs.slides)
        s_title = slides[IDX_TITLE]
        s_toc = slides[IDX_TOC]
        s_summary = slides[IDX_SUMMARY]
        s_div1 = slides[IDX_DIVIDER_1]
        s_intro = slides[IDX_INTRO]
        s_workflow = slides[IDX_WORKFLOW]
        s_sample = slides[IDX_SAMPLE]
        s_div2 = slides[IDX_DIVIDER_2]
        s_results = slides[IDX_RESULTS]
        s_div3 = slides[IDX_DIVIDER_3]
        s_conclusion = slides[IDX_CONCLUSION]
        # slides[IDX_END] (blank closing slide) and slides[IDX_COLOR_GUIDE]
        # (page 1) are left completely untouched.

        self._fill_title_page(s_title, data.title_page)
        self._fill_toc(s_toc, data.chapters)
        self._fill_summary(s_summary, data.summary)
        self._fill_divider(s_div1, data.divider_intro)
        self._fill_section_title_only(s_intro)  # "Introduction" heading, static
        self._fill_workflow(s_workflow, s_sample, data.workflow_content)
        self._fill_sample_info(s_sample, data.sample_info)

        if data.divider_results is not None:
            self._fill_divider(s_div2, data.divider_results)
        else:
            su.delete_slide_obj(self.prs, s_div2)

        self._fill_results(s_results, data.results)

        if data.divider_conclusion is not None:
            self._fill_divider(s_div3, data.divider_conclusion)
        else:
            su.delete_slide_obj(self.prs, s_div3)

        if data.conclusion is not None:
            self._fill_conclusion(s_conclusion, data.conclusion)
        else:
            su.delete_slide_obj(self.prs, s_conclusion)

        self.prs.save(output_path)
        return output_path

    # ------------------------------------------------------------------
    # Page 2 -- Title page
    # ------------------------------------------------------------------
    def _fill_title_page(self, slide, title_data: TitlePageData):
        title_shape = su.find_shape_by_text(slide, "[Title]")
        if title_shape:
            su.set_single_line_text(title_shape, title_data.title)

        subtitle_shape = su.find_shape_by_text(slide, "[subtitle]")
        if subtitle_shape:
            su.set_single_line_text(subtitle_shape, title_data.subtitle)

        byline_shape = su.find_shape_by_text(slide, "[writer] ｜[team] ｜[date]")
        if byline_shape is None:
            # fall back to a prefix search in case of stray whitespace
            matches = su.find_shapes_starting_with(slide, "[writer]")
            byline_shape = matches[0] if matches else None
        if byline_shape:
            byline_text = f"{title_data.writer} ｜ {title_data.team} ｜ {title_data.date}"
            su.set_single_line_text(byline_shape, byline_text)

    # ------------------------------------------------------------------
    # Page 3 -- Table of contents (up to 4 chapter slots)
    # ------------------------------------------------------------------
    def _fill_toc(self, slide, chapters):
        if len(chapters) > TOC_MAX_CHAPTERS:
            raise ValueError(
                f"TOC layout supports at most {TOC_MAX_CHAPTERS} chapters, got {len(chapters)}"
            )
        for slot in range(1, TOC_MAX_CHAPTERS + 1):
            group = su.find_group_containing(slide, f"CHAPTER {slot}")
            if group is None:
                continue
            if slot <= len(chapters):
                label_shape = su.find_shape_by_text(group, "제목을 입력해 주십시오")
                # find_shape_by_text expects a slide-like object with .shapes
                if label_shape is None:
                    for sub in su.iter_shapes_recursive(group.shapes):
                        if sub.has_text_frame and sub.text_frame.text.strip() == "제목을 입력해 주십시오":
                            label_shape = sub
                            break
                if label_shape is not None:
                    su.set_single_line_text(label_shape, chapters[slot - 1].name)
            else:
                # fewer chapters than slots -- remove the unused quadrant entirely
                group._element.getparent().remove(group._element)

    # ------------------------------------------------------------------
    # Page 4 -- Summary (Purpose / Results / Conclusion / Further Study)
    # ------------------------------------------------------------------
    def _fill_summary(self, slide, summary: SummaryData):
        quadrants = {
            "Purpose": summary.purpose,
            "Results": summary.results,
            "Conclusion": summary.conclusion,
            "Further Study": summary.further_study,
        }
        for label, items in quadrants.items():
            group = su.find_group_containing(slide, label)
            if group is None:
                continue
            content_shape = None
            for sub in su.iter_shapes_recursive(group.shapes):
                if sub.has_text_frame and sub.text_frame.text.strip().startswith("내용을 입력해 주십시오"):
                    content_shape = sub
                    break
            if content_shape is not None:
                su.fill_bullets(content_shape, items or [""], bulleted=False)

    # ------------------------------------------------------------------
    # Chapter divider slides ("CHAPTER 0N" / name)
    # ------------------------------------------------------------------
    def _fill_divider(self, slide, divider: DividerData):
        number_shape = su.find_shape_by_text(slide, "[Number]")
        if number_shape:
            su.set_single_line_text(number_shape, divider.number)
        name_shape = su.find_shape_by_text(slide, "[Chapter Name]")
        if name_shape:
            su.set_single_line_text(name_shape, divider.name)

    def _fill_section_title_only(self, slide):
        # "Introduction" / "Workflow" section-header slides carry no other
        # placeholders in the base template -- nothing to fill.
        pass

    # ------------------------------------------------------------------
    # Workflow: optional content slide, reusing the Sample-Info layout
    # ------------------------------------------------------------------
    def _fill_workflow(self, workflow_title_slide, sample_layout_slide, content):
        if content is None:
            return
        new_slide = su.duplicate_slide(self.prs, self._index_of(sample_layout_slide))
        su.move_slide_after(self.prs, new_slide, workflow_title_slide)

        title_ph = new_slide.shapes.title
        if title_ph is not None:
            su.set_single_line_text(title_ph, "Workflow")

        content_shape = su.find_shape_by_text(new_slide, "[CONTENTS]")
        if content_shape is not None:
            apply_content_block(new_slide, content_shape, content)

    # ------------------------------------------------------------------
    # Sample / Data information
    # ------------------------------------------------------------------
    def _fill_sample_info(self, slide, sample_info):
        if sample_info is None:
            su.delete_slide_obj(self.prs, slide)
            return
        content_shape = su.find_shape_by_text(slide, "[CONTENTS]")
        if content_shape is not None:
            apply_content_block(slide, content_shape, sample_info.content)

    # ------------------------------------------------------------------
    # Results (repeatable)
    # ------------------------------------------------------------------
    def _fill_results(self, template_slide, results):
        if not results:
            su.delete_slide_obj(self.prs, template_slide)
            return

        # Duplicate the *pristine* template (still has [TITLE]/[CONTENTS]
        # placeholders untouched) for every extra item BEFORE filling any of
        # them in. Filling first and duplicating after would copy the
        # already-replaced text, so later items would silently fail to find
        # their placeholders.
        targets = [template_slide]
        anchor = template_slide
        for _ in results[1:]:
            new_slide = su.duplicate_slide(self.prs, self._index_of(template_slide))
            su.move_slide_after(self.prs, new_slide, anchor)
            anchor = new_slide
            targets.append(new_slide)

        for target, item in zip(targets, results):
            title_ph = target.shapes.title
            if title_ph is not None:
                su.replace_text_token(title_ph, "[TITLE]", item.title)

            label_shape = su.find_shape_by_text(target, "[TITLE]")
            if label_shape is not None:
                su.set_single_line_text(label_shape, item.title)

            content_shape = su.find_shape_by_text(target, "[CONTENTS]")
            if content_shape is not None:
                apply_content_block(target, content_shape, item.content)

    # ------------------------------------------------------------------
    # Conclusion
    # ------------------------------------------------------------------
    def _fill_conclusion(self, slide, conclusion):
        title_shape = su.find_shape_by_text(slide, "[TITLE]")
        if title_shape is not None:
            su.set_single_line_text(title_shape, conclusion.title)

        content_shapes = [
            s for s in su.iter_shapes_recursive(slide.shapes)
            if s.has_text_frame and s.text_frame.text.strip() == "[CONTENTS]"
        ]
        content_shapes.sort(key=lambda s: s.top)

        if content_shapes and conclusion.content_1 is not None:
            apply_content_block(slide, content_shapes[0], conclusion.content_1)
        if len(content_shapes) > 1 and conclusion.content_2 is not None:
            apply_content_block(slide, content_shapes[1], conclusion.content_2)

    # ------------------------------------------------------------------
    def _index_of(self, slide) -> int:
        for i, s in enumerate(self.prs.slides):
            if s.slide_id == slide.slide_id:
                return i
        raise ValueError("slide not found in presentation")
